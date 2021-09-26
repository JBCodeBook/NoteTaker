import glob
import os.path
from pathlib import Path
from pptx import Presentation


def takeNotes(f, noteFile):
    prs = Presentation(f)
    text_runs = []
    indexOf = []
    picList = []
    isTitle = True
    id = 0

    for index, slide in enumerate(prs.slides):
        isFirst = True

        for shape in slide.shapes:

            if not shape.has_text_frame:

                if ("Picture" in shape.name):
                    lecture_name = os.path.basename(f)[0:-5]
                    lecture_name = "_".join(lecture_name.split("_", 2)[:2])
                    Path("pics/" + lecture_name).mkdir(parents=True, exist_ok=True)
                    path_to_pictures = 'pics/' + lecture_name + '/mypic' + str(id) + '.jpg'

                    with open(path_to_pictures, 'wb') as picture:
                         picture.write(shape.image.blob)
                         tup = (index, path_to_pictures)
                         picList.append(tup)
                         id = id + 1



                continue

            for paragraph in shape.text_frame.paragraphs:




                for run in paragraph.runs:
                    if len(run.text) > 0:




                        if isTitle:
                            text_runs.append("*** " + run.text)
                            isTitle = False

                        if isFirst:
                            text_runs.append("**** " + run.text)
                            isFirst = False

                            break
                        text_runs.append("- " + run.text)

    text_runs.append("\n")

    for i in picList:
        print(i)

    # for notes in text_runs:
    #     noteFile.write(notes + "\n")


def main():
    title = "#+TITLE:"
    author = "#+AUTHOR:"
    opt = "#+OPTIONS: toc:nil num:nil"
    style = "#+SETUPFILE: C:/Users/j2017/OneDrive/School/Org-Mode/org-html-themes-master/org-html-themes-master/org/theme-bigblow-local.setup"

    noteList = glob.glob("C:/Users/j2017/OneDrive/School/pysc\\Lectures_powerpoint/*.pptx")

    with open("PyscNotes" + ".org", "a+") as noteFile:
        noteFile.write(title + "pysch" + "\n")
        noteFile.write(author + "John B" + "\n")
        noteFile.write(opt + "\n")
        noteFile.write(style + "\n\n")

        for index, files in enumerate(noteList):

            if index == 0:
                noteFile.write("* " + os.path.basename(files)[0:-5] + "\n")
                takeNotes(files, noteFile)

            if index > 0:
                current_file = os.path.basename(files)[0:-5]
                current_file = "_".join(current_file.split("_", 2)[:2])

                previous_file = os.path.basename(noteList[index - 1])[0:-5]
                previous_file = "_".join(previous_file.split("_", 2)[:2])

                if current_file == previous_file:
                    noteFile.write("** " + os.path.basename(files)[0:-5] + "\n")
                    takeNotes(files, noteFile)

                else:
                    noteFile.write("* " + os.path.basename(files)[0:-5] + "\n")
                    takeNotes(files, noteFile)

    print("Closing File \n\n\n\n")
    noteFile.close()


if __name__ == "__main__":
    main()

# TODO you need to bring back the part of the loop that extracts the actual text from each heading in the paragraph loop
